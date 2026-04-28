#!/usr/bin/env python3
from __future__ import annotations

import re
import shutil
from pathlib import Path

import cv2
import easyocr
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.util import Pt


ROOT = Path(__file__).resolve().parent
SOURCE_IMAGES = sorted(ROOT.glob("source300-slide-*.jpg"))
TARGET_PPT = ROOT / "Towards a Science of Scaling Agent Systems_中文解读.pptx"
BACKUP_PPT = ROOT / "Towards a Science of Scaling Agent Systems_中文解读.before_editable_kaiti_backup.pptx"
WORK_DIR = ROOT / ".kaiti_editable_build"
BG_DIR = WORK_DIR / "backgrounds"
PREVIEW_GRID = WORK_DIR / "editable-preview-grid.jpg"


REPLACEMENTS = {
    "硬按": "硬核",
    "深餍": "深度",
    "井行": "并行",
    "准确孪": "准确率",
    "趴服从度": "服从度",
    "交叉~证": "交叉验证",
    "错误放大宰": "错误放大率",
    "性能幂跌": "性能暴跌",
}

NOISE_TEXTS = {
    "畿",
    "C裟",
    "[贝",
    "GC二",
    "贪",
}


def normalize_text(text: str) -> str:
    t = text.strip()
    for a, b in REPLACEMENTS.items():
        t = t.replace(a, b)
    return t


def is_useful_text(text: str, conf: float, w: int, h: int) -> bool:
    if not re.search(r"[\u4e00-\u9fff]", text):
        return False
    if text in NOISE_TEXTS:
        return False

    area = w * h
    cjk_count = len(re.findall(r"[\u4e00-\u9fff]", text))
    cjk_ratio = cjk_count / max(1, len(text))

    if len(text) <= 1 and conf < 0.95:
        return False
    if h > w * 1.35:
        return False
    if area < 1200 and conf < 0.70:
        return False
    if conf < 0.18 and cjk_count < 6:
        return False
    if cjk_ratio < 0.22 and conf < 0.50:
        return False
    return True


def estimate_text_rgb(img: np.ndarray, x: int, y: int, w: int, h: int) -> tuple[int, int, int]:
    roi = img[max(0, y) : min(img.shape[0], y + h), max(0, x) : min(img.shape[1], x + w)]
    if roi.size == 0:
        return (20, 30, 60)

    gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
    k = max(20, int(gray.size * 0.06))
    idx = np.argpartition(gray.reshape(-1), k - 1)[:k]
    col = np.median(roi.reshape(-1, 3)[idx], axis=0).astype(int)  # BGR
    return (int(col[2]), int(col[1]), int(col[0]))  # RGB


def remove_text_from_background(img: np.ndarray, items: list[dict]) -> np.ndarray:
    cleaned = img.copy()
    h, w = cleaned.shape[:2]
    for item in items:
        pts = item["pts"]
        x, y, bw, bh = item["bbox"]

        mask = np.zeros((h, w), dtype=np.uint8)
        cv2.fillPoly(mask, [pts], 255)

        x1, y1 = max(0, x - 6), max(0, y - 6)
        x2, y2 = min(w, x + bw + 6), min(h, y + bh + 6)
        ring = np.zeros((h, w), dtype=np.uint8)
        cv2.rectangle(ring, (x1, y1), (x2, y2), 255, -1)
        ring = cv2.bitwise_and(ring, cv2.bitwise_not(mask))
        ring_pixels = cleaned[ring == 255]
        if ring_pixels.size:
            bg = np.median(ring_pixels, axis=0).astype(np.uint8)  # BGR
        else:
            bg = np.array([255, 255, 255], dtype=np.uint8)
        cleaned[mask == 255] = bg
    return cleaned


def parse_slide(
    image_path: Path, reader: easyocr.Reader
) -> tuple[np.ndarray, list[dict], int, int]:
    img = cv2.imread(str(image_path))
    if img is None:
        raise RuntimeError(f"Cannot load image: {image_path}")
    h, w = img.shape[:2]

    raw = reader.readtext(str(image_path), detail=1, paragraph=False)
    items: list[dict] = []
    for box, text, conf in raw:
        text = normalize_text(text)
        pts = np.array(box, dtype=np.int32)
        x, y, bw, bh = cv2.boundingRect(pts)
        if not is_useful_text(text, conf, bw, bh):
            continue
        rgb = estimate_text_rgb(img, x, y, bw, bh)
        items.append(
            {
                "pts": pts,
                "text": text,
                "conf": conf,
                "bbox": (x, y, bw, bh),
                "rgb": rgb,
            }
        )

    cleaned = remove_text_from_background(img, items)
    return cleaned, items, w, h


def create_preview_grid(backgrounds: list[Path], out_path: Path) -> None:
    thumbs = []
    for p in backgrounds:
        im = Image.open(p).convert("RGB")
        thumbs.append(im.resize((480, 270), Image.Resampling.LANCZOS))
    cols = 3
    rows = (len(thumbs) + cols - 1) // cols
    canvas = Image.new("RGB", (20 + cols * 520, 20 + rows * 330), "white")
    for i, thumb in enumerate(thumbs):
        r, c = divmod(i, cols)
        x = 20 + c * 520
        y = 20 + r * 330
        canvas.paste(thumb, (x, y))
    canvas.save(out_path, quality=92)


def build_editable_ppt(slide_data: list[dict], out_path: Path) -> None:
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = 12192000
    prs.slide_height = 6858000
    blank = prs.slide_layouts[6]

    for data in slide_data:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(data["bg_path"]), 0, 0, width=prs.slide_width, height=prs.slide_height
        )

        px_to_emu_x = prs.slide_width / data["img_w"]
        px_to_emu_y = prs.slide_height / data["img_h"]
        pt_per_px = 540.0 / data["img_h"]  # 7.5 in * 72 pt

        for item in data["items"]:
            x, y, w, h = item["bbox"]
            left = int(x * px_to_emu_x)
            top = int(y * px_to_emu_y)
            width = max(int(w * px_to_emu_x), 8000)
            height = max(int(h * px_to_emu_y), 8000)
            textbox = slide.shapes.add_textbox(left, top, width, height)

            tf = textbox.text_frame
            tf.clear()
            tf.margin_left = 0
            tf.margin_right = 0
            tf.margin_top = 0
            tf.margin_bottom = 0
            tf.word_wrap = True
            tf.auto_size = MSO_AUTO_SIZE.NONE
            tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE

            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = item["text"]

            # Slightly smaller than bbox height to keep inside line box.
            size_pt = max(9, min(44, h * pt_per_px * 0.80))
            run.font.size = Pt(size_pt)
            run.font.name = "STKaiti"
            run.font.color.rgb = RGBColor(*item["rgb"])

    prs.save(out_path)


def main() -> None:
    if not SOURCE_IMAGES:
        raise FileNotFoundError("No source300-slide-*.jpg files found.")

    WORK_DIR.mkdir(exist_ok=True)
    BG_DIR.mkdir(exist_ok=True)

    reader = easyocr.Reader(["ch_sim", "en"], gpu=False, verbose=False)
    slide_data: list[dict] = []
    bg_paths: list[Path] = []

    for img_path in SOURCE_IMAGES:
        cleaned, items, w, h = parse_slide(img_path, reader)
        bg_path = BG_DIR / img_path.name.replace("source300-", "bg-")
        cv2.imwrite(str(bg_path), cleaned, [int(cv2.IMWRITE_JPEG_QUALITY), 92])
        bg_paths.append(bg_path)
        slide_data.append(
            {"bg_path": bg_path, "items": items, "img_w": w, "img_h": h, "src": img_path}
        )
        print(f"{img_path.name}: editable text blocks={len(items)}")

    create_preview_grid(bg_paths, PREVIEW_GRID)

    if TARGET_PPT.exists():
        shutil.copy2(TARGET_PPT, BACKUP_PPT)
        print(f"Backup created: {BACKUP_PPT.name}")

    build_editable_ppt(slide_data, TARGET_PPT)
    print(f"Editable PPT updated: {TARGET_PPT}")
    print(f"Preview grid: {PREVIEW_GRID}")


if __name__ == "__main__":
    main()
