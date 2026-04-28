#!/usr/bin/env python3
from __future__ import annotations

import re
import shutil
from pathlib import Path

import cv2
import easyocr
import numpy as np
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches


ROOT = Path(__file__).resolve().parent
SOURCE_IMAGES_GLOB = "source300-slide-*.jpg"
TARGET_PPT = ROOT / "Towards a Science of Scaling Agent Systems_中文解读.pptx"
TARGET_BACKUP = ROOT / "Towards a Science of Scaling Agent Systems_中文解读.before_kaiti_backup.pptx"
WORK_DIR = ROOT / ".kaiti_build"
PROCESSED_DIR = WORK_DIR / "processed_images"
GRID_PATH = WORK_DIR / "processed-grid.jpg"


def find_kaiti_font() -> Path:
    asset_fonts = sorted(
        Path("/System/Library/AssetsV2").glob(
            "com_apple_MobileAsset_Font8/*.asset/AssetData/Kaiti.ttc"
        )
    )
    if asset_fonts:
        return asset_fonts[-1]

    fallback = Path(
        "/System/Library/PrivateFrameworks/FontServices.framework/Versions/A/"
        "Resources/Fonts/Subsets/Kaiti.ttc"
    )
    if fallback.exists():
        return fallback
    raise FileNotFoundError("Cannot find Kaiti.ttc on this system.")


def normalize_text(text: str) -> str:
    replacements = {
        "硬按": "硬核",
        "深餍": "深度",
        "井行": "并行",
        "准确孪": "准确率",
        "趴服从度": "服从度",
        "交叉~证": "交叉验证",
        "效率 (0%)": "效率 (E)",
    }
    for src, dst in replacements.items():
        text = text.replace(src, dst)
    return text


def should_keep_item(text: str, conf: float, bw: int, bh: int) -> bool:
    cjk_count = len(re.findall(r"[\u4e00-\u9fff]", text))
    if cjk_count == 0:
        return False

    area = bw * bh
    cjk_ratio = cjk_count / max(1, len(text))

    if bh > bw * 1.35:
        return False
    if len(text.strip()) <= 1 and conf < 0.9:
        return False
    if conf < 0.18 and cjk_count < 6:
        return False
    if cjk_ratio < 0.25 and conf < 0.5:
        return False
    if area < 1200 and conf < 0.65:
        return False
    return True


def replace_chinese_text(
    image_path: Path, reader: easyocr.Reader, font_path: Path
) -> tuple[Image.Image, int]:
    bgr = cv2.imread(str(image_path))
    if bgr is None:
        raise RuntimeError(f"Failed to load image: {image_path}")

    detected = reader.readtext(str(image_path), detail=1, paragraph=False)
    items: list[tuple[np.ndarray, str, float, tuple[int, int, int, int]]] = []
    h, w = bgr.shape[:2]

    for box, raw_text, conf in detected:
        if not re.search(r"[\u4e00-\u9fff]", raw_text):
            continue
        pts = np.array(box, dtype=np.int32)
        x, y, bw, bh = cv2.boundingRect(pts)
        text = normalize_text(raw_text)
        if not should_keep_item(text, conf, bw, bh):
            continue
        items.append((pts, text, conf, (x, y, bw, bh)))

    cleaned = bgr.copy()
    for pts, _, _, (x, y, bw, bh) in items:
        mask = np.zeros((h, w), dtype=np.uint8)
        cv2.fillPoly(mask, [pts], 255)

        x1, y1 = max(0, x - 6), max(0, y - 6)
        x2, y2 = min(w, x + bw + 6), min(h, y + bh + 6)

        ring = np.zeros((h, w), dtype=np.uint8)
        cv2.rectangle(ring, (x1, y1), (x2, y2), 255, -1)
        ring = cv2.bitwise_and(ring, cv2.bitwise_not(mask))
        ring_pixels = cleaned[ring == 255]

        if ring_pixels.size:
            bg_color = np.median(ring_pixels, axis=0).astype(np.uint8)
        else:
            bg_color = np.array([255, 255, 255], dtype=np.uint8)
        cleaned[mask == 255] = bg_color

    pil = Image.fromarray(cv2.cvtColor(cleaned, cv2.COLOR_BGR2RGB))
    draw = ImageDraw.Draw(pil)

    for _, text, _, (x, y, bw, bh) in items:
        size = max(12, int(bh * 0.78))
        font = ImageFont.truetype(str(font_path), size=size, index=0)
        box = draw.textbbox((0, 0), text, font=font)
        tw, th = box[2] - box[0], box[3] - box[1]
        while (tw > bw * 1.03 or th > bh * 1.12) and size > 9:
            size -= 1
            font = ImageFont.truetype(str(font_path), size=size, index=0)
            box = draw.textbbox((0, 0), text, font=font)
            tw, th = box[2] - box[0], box[3] - box[1]

        roi = bgr[max(0, y) : min(h, y + bh), max(0, x) : min(w, x + bw)]
        if roi.size:
            gray = cv2.cvtColor(roi, cv2.COLOR_BGR2GRAY)
            k = max(20, int(gray.size * 0.06))
            idx = np.argpartition(gray.reshape(-1), k - 1)[:k]
            col = np.median(roi.reshape(-1, 3)[idx], axis=0).astype(int)
            fill = (int(col[2]), int(col[1]), int(col[0]))
        else:
            fill = (20, 30, 60)

        draw.text((x, y), text, font=font, fill=fill)

    # Keep file size under control while retaining good visual quality.
    if pil.width > 3200:
        pil = pil.resize((pil.width // 2, pil.height // 2), Image.Resampling.LANCZOS)

    return pil, len(items)


def make_grid(image_paths: list[Path], out_path: Path, cols: int = 3) -> None:
    thumbs = []
    for p in image_paths:
        img = Image.open(p).convert("RGB")
        thumbs.append(img.resize((480, 270), Image.Resampling.LANCZOS))

    rows = (len(thumbs) + cols - 1) // cols
    canvas = Image.new("RGB", (cols * 520 + 20, rows * 330 + 20), "white")
    draw = ImageDraw.Draw(canvas)
    for idx, thumb in enumerate(thumbs):
        r, c = divmod(idx, cols)
        x = 20 + c * 520
        y = 20 + r * 330
        canvas.paste(thumb, (x, y))
        draw.text((x, y + 278), image_paths[idx].stem, fill="black")
    canvas.save(out_path, quality=92)


def build_ppt(image_paths: list[Path], out_ppt: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for img_path in image_paths:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(img_path), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    prs.save(out_ppt)


def main() -> None:
    source_images = sorted(ROOT.glob(SOURCE_IMAGES_GLOB))
    if not source_images:
        raise FileNotFoundError(
            f"No source images found by pattern: {SOURCE_IMAGES_GLOB}"
        )

    WORK_DIR.mkdir(exist_ok=True)
    PROCESSED_DIR.mkdir(exist_ok=True)

    font_path = find_kaiti_font()
    reader = easyocr.Reader(["ch_sim", "en"], gpu=False, verbose=False)

    processed: list[Path] = []
    for src in source_images:
        out_img = PROCESSED_DIR / src.name.replace("source300-", "kaiti-")
        pil, n = replace_chinese_text(src, reader, font_path)
        pil.save(out_img, quality=92)
        processed.append(out_img)
        print(f"{src.name}: replaced {n} text block(s)")

    make_grid(processed, GRID_PATH)

    if TARGET_PPT.exists() and not TARGET_BACKUP.exists():
        shutil.copy2(TARGET_PPT, TARGET_BACKUP)
        print(f"Backup created: {TARGET_BACKUP.name}")

    build_ppt(processed, TARGET_PPT)
    print(f"Updated PPT: {TARGET_PPT}")
    print(f"Preview grid: {GRID_PATH}")


if __name__ == "__main__":
    main()
