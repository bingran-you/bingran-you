from __future__ import annotations

from pathlib import Path
from typing import Iterable

from PIL import Image, ImageDraw, ImageFont, ImageFilter
from pptx import Presentation


ROOT = Path(
    "/Users/bingran_you/Library/Mobile Documents/com~apple~CloudDocs/Downloads/huaying/LFM2-VL"
)
TEMPLATE = ROOT / "PPT模板.pptx"
OUTPUT = ROOT / "LFM2-VL_图解讲解版.pptx"
ASSET_DIR = ROOT / "generated_assets"
WEB_ASSET_DIR = ROOT / "LFM2-VL_ Efficient Vision-Language Models _ Liquid AI_files"

FONT_CANDIDATES = [
    "/System/Library/Fonts/Hiragino Sans GB.ttc",
    "/System/Library/Fonts/STHeiti Medium.ttc",
    "/System/Library/Fonts/STHeiti Light.ttc",
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    "/Library/Fonts/Arial Unicode.ttf",
]

COLORS = {
    "bg": "#FFFFFF",
    "card": "#F5F5F5",
    "border": "#D8D8D8",
    "title": "#222222",
    "body": "#4A4A4A",
    "accent": "#C1372A",
    "muted": "#7A7A7A",
    "black": "#111111",
}


def font_path() -> str:
    for candidate in FONT_CANDIDATES:
        if Path(candidate).exists():
            return candidate
    raise FileNotFoundError("No usable CJK font found")


FONT = font_path()


def load_font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont:
    return ImageFont.truetype(FONT, size=size)


def wrap_text(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    if not text:
        return [""]
    lines: list[str] = []
    current = ""
    for ch in text:
        candidate = current + ch
        if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
            current = candidate
        else:
            if current:
                lines.append(current)
            current = ch
    if current:
        lines.append(current)
    return lines or [""]


def draw_paragraph(
    draw: ImageDraw.ImageDraw,
    x: int,
    y: int,
    text: str,
    font: ImageFont.FreeTypeFont,
    fill: str,
    max_width: int,
    line_gap: int = 10,
) -> int:
    lines = wrap_text(draw, text, font, max_width)
    line_height = draw.textbbox((0, 0), "测试Ag", font=font)[3] + line_gap
    for i, line in enumerate(lines):
        draw.text((x, y + i * line_height), line, font=font, fill=fill)
    return y + len(lines) * line_height


def rounded_box(draw: ImageDraw.ImageDraw, xy: tuple[int, int, int, int], radius: int = 28, fill: str = COLORS["card"], outline: str = COLORS["border"], width: int = 2) -> None:
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def create_edge_intro() -> Path:
    out = ASSET_DIR / "edge_intro.png"
    img = Image.new("RGB", (2000, 760), COLORS["bg"])
    draw = ImageDraw.Draw(img)

    title_font = load_font(52)
    body_font = load_font(28)
    small_font = load_font(24)

    cards = [
        ("低延迟", "终端问答与拍照理解需要快速响应，不能把多模态推理全部留给云端。"),
        ("低资源占用", "显存、内存和能耗都受限，传统大体量 VLM 很难直接落在手机和边缘盒子。"),
        ("可调推理成本", "不同图像分辨率和任务复杂度，需要可以在线调节的 token 预算。"),
    ]

    card_w = 570
    gap = 45
    start_x = 95
    y1, y2 = 120, 640

    draw.text((95, 40), "边缘级 VLM 的三重约束", font=title_font, fill=COLORS["title"])
    draw.text((98, 90), "LFM2-VL 试图同时解决速度、内存和灵活推理三个问题。", font=small_font, fill=COLORS["muted"])

    for idx, (label, text) in enumerate(cards):
        x1 = start_x + idx * (card_w + gap)
        x2 = x1 + card_w
        rounded_box(draw, (x1, y1, x2, y2))
        draw.rounded_rectangle((x1 + 24, y1 + 24, x1 + 78, y1 + 78), radius=14, fill=COLORS["accent"])
        draw.text((x1 + 105, y1 + 24), label, font=title_font, fill=COLORS["black"])
        draw_paragraph(draw, x1 + 38, y1 + 120, text, body_font, COLORS["body"], card_w - 76, line_gap=14)

    img.save(out)
    return out


def create_token_flow() -> Path:
    out = ASSET_DIR / "token_flow.png"
    img = Image.new("RGB", (2000, 760), COLORS["bg"])
    draw = ImageDraw.Draw(img)
    step_font = load_font(34)
    body_font = load_font(24)

    steps = [
        ("1 输入", "文本提示 + 原图"),
        ("2 切片", "小图原生输入\n大图切成 512x512 patch"),
        ("3 编码", "SigLIP2 NaFlex\n提取视觉 token"),
        ("4 压缩", "PixelUnshuffle +\n2-layer MLP Connector"),
        ("5 推理", "与文本 token 融合\n交给 LFM2 输出"),
    ]

    box_w, box_h = 290, 170
    start_x, y = 105, 150
    gap = 35
    for idx, (label, text) in enumerate(steps):
        x1 = start_x + idx * (box_w + gap)
        x2 = x1 + box_w
        rounded_box(draw, (x1, y, x2, y + box_h), radius=24)
        draw.rounded_rectangle((x1 + 18, y + 18, x1 + 72, y + 72), radius=12, fill=COLORS["accent"])
        draw.text((x1 + 88, y + 18), label, font=step_font, fill=COLORS["black"])
        current_y = y + 92
        for line in text.split("\n"):
            draw.text((x1 + 26, current_y), line, font=body_font, fill=COLORS["body"])
            current_y += 34
        if idx < len(steps) - 1:
            mid_y = y + box_h // 2
            x_arrow = x2 + 16
            draw.line((x_arrow, mid_y, x_arrow + 18, mid_y), fill=COLORS["accent"], width=6)
            draw.polygon(
                [(x_arrow + 18, mid_y - 10), (x_arrow + 18, mid_y + 10), (x_arrow + 36, mid_y)],
                fill=COLORS["accent"],
            )

    footer = "论文例子：256x384 -> 96 tokens，384x680 -> 240 tokens，1000x3000 -> 1020 tokens。"
    draw.text((520, 645), footer, font=body_font, fill=COLORS["body"])
    img.save(out)
    return out


def create_benchmark_table() -> Path:
    out = ASSET_DIR / "benchmark_scorecard.png"
    img = Image.new("RGB", (2000, 760), COLORS["bg"])
    draw = ImageDraw.Draw(img)
    head_font = load_font(28)
    body_font = load_font(25)

    headers = ["任务", "LFM2-VL-1.6B", "LFM2-VL-450M", "对比 2B 级模型"]
    rows = [
        ("RealWorldQA", "65.23", "52.29", "65.10"),
        ("InfoVQA (Val)", "58.68", "46.51", "66.10"),
        ("OCRBench", "742", "655", "831"),
        ("MathVista", "51.10", "44.70", "57.60"),
        ("SEEDBench_IMG", "71.97", "63.50", "75.00"),
    ]

    table_x, table_y = 140, 90
    col_w = [430, 350, 350, 430]
    row_h = 88

    x = table_x
    for idx, header in enumerate(headers):
        rounded_box(draw, (x, table_y, x + col_w[idx], table_y + row_h), radius=14, fill="#F0F0F0")
        draw.text((x + 20, table_y + 24), header, font=head_font, fill=COLORS["black"])
        x += col_w[idx] + 8

    for ridx, row in enumerate(rows):
        y = table_y + row_h + 12 + ridx * (row_h + 8)
        x = table_x
        row_fill = "#FFFFFF" if ridx % 2 == 0 else "#FAFAFA"
        for cidx, value in enumerate(row):
            rounded_box(draw, (x, y, x + col_w[cidx], y + row_h), radius=14, fill=row_fill)
            fill = COLORS["accent"] if cidx == 1 else COLORS["body"]
            draw.text((x + 20, y + 24), value, font=body_font, fill=fill)
            x += col_w[cidx] + 8

    img.save(out)
    return out


def create_variant_compare() -> Path:
    out = ASSET_DIR / "variant_compare.png"
    img = Image.new("RGB", (1800, 650), COLORS["bg"])
    draw = ImageDraw.Draw(img)
    head_font = load_font(42)
    body_font = load_font(28)
    tag_font = load_font(34)

    cards = [
        (
            "LFM2-VL-450M",
            [
                "更轻量",
                "手机 / 穿戴 / 边缘盒子",
            ],
        ),
        (
            "LFM2-VL-1.6B",
            [
                "更强细节",
                "笔记本 / 单 GPU / 本地工作站",
            ],
        ),
    ]

    card_w, card_h = 760, 360
    positions = [(70, 145), (970, 145)]
    for (title, bullets), (x, y) in zip(cards, positions):
        rounded_box(draw, (x, y, x + card_w, y + card_h), radius=28)
        draw.rounded_rectangle((x + 28, y + 28, x + 90, y + 90), radius=14, fill=COLORS["accent"])
        draw.text((x + 112, y + 30), title, font=head_font, fill=COLORS["black"])
        draw.text((x + 42, y + 135), bullets[0], font=tag_font, fill=COLORS["accent"])
        draw_paragraph(draw, x + 42, y + 205, bullets[1], body_font, COLORS["body"], card_w - 84, line_gap=10)

    img.save(out)
    return out


def create_cover_banner() -> Path:
    out = ASSET_DIR / "cover_banner.png"
    canvas = Image.new("RGB", (2400, 780), COLORS["bg"])
    source = Image.open(WEB_ASSET_DIR / "689b92d8747e3502dcd4db50_Artboard 1_1.jpg").convert("RGB")
    photo = source.crop((620, 55, 1380, 700)).resize((2400, 780)).filter(ImageFilter.GaussianBlur(1.2))
    canvas.paste(photo, (0, 0))
    canvas.save(out)
    return out


def create_demo_showcase() -> Path:
    out = ASSET_DIR / "demo_showcase.png"
    bg = Image.new("RGB", (2000, 760), COLORS["bg"])
    draw = ImageDraw.Draw(bg)
    head_font = load_font(36)
    body_font = load_font(28)
    quote_font = load_font(34)

    source = Image.open(WEB_ASSET_DIR / "689b92d8747e3502dcd4db50_Artboard 1_1.jpg").convert("RGB")
    photo = source.crop((620, 55, 1380, 700)).resize((860, 640))
    bg.paste(photo, (90, 60))

    rounded_box(draw, (1040, 120, 1890, 560), radius=24)
    draw.text((1085, 160), "Figure 1 demo 的直观信息", font=head_font, fill=COLORS["black"])
    draw_paragraph(
        draw,
        1085,
        230,
        "模型不仅回答了“这是海滩”，还识别出前景中的牛，并把场景组织成一段完整自然语言描述。",
        quote_font,
        COLORS["body"],
        730,
        line_gap=14,
    )
    draw_paragraph(
        draw,
        1085,
        410,
        "这说明 450M 版本已经具备基础视觉 grounding、场景概括和文字生成的一体化能力。",
        body_font,
        COLORS["body"],
        730,
        line_gap=12,
    )
    bg.save(out)
    return out


def ensure_assets() -> dict[str, Path]:
    ASSET_DIR.mkdir(exist_ok=True)
    return {
        "edge_intro": create_edge_intro(),
        "token_flow": create_token_flow(),
        "benchmark": create_benchmark_table(),
        "variant_compare": create_variant_compare(),
        "cover_banner": create_cover_banner(),
        "demo_showcase": create_demo_showcase(),
        "architecture": WEB_ASSET_DIR / "689bafd84898860cd0c654db_Architecture (5).png",
        "demo": WEB_ASSET_DIR / "689b92d8747e3502dcd4db50_Artboard 1_1.jpg",
        "speed": WEB_ASSET_DIR / "689b3eef2ae10f1ac5e8c338_LFM2-VL - Vision-Language Models_ Processing Time Comparison (4).png",
        "memory": WEB_ASSET_DIR / "689b7df6b9722eceb23cc309_LFM2-VL - Vision-Language Models_ Memory Footprint (4).png",
    }


def get_shape(container, path: Iterable[int]):
    shape = None
    shapes = container.shapes if hasattr(container, "shapes") else container
    for idx in path:
        shape = shapes[idx - 1]
        shapes = shape.shapes if hasattr(shape, "shapes") else None
    return shape


def set_paragraphs(shape, lines: list[str]) -> None:
    tf = shape.text_frame
    while len(tf.paragraphs) < len(lines):
        tf.add_paragraph()
    for idx, paragraph in enumerate(tf.paragraphs):
        text = lines[idx] if idx < len(lines) else ""
        if paragraph.runs:
            paragraph.runs[0].text = text
            for run in paragraph.runs[1:]:
                run.text = ""
        else:
            paragraph.add_run().text = text


def replace_picture(slide, shape_path: Iterable[int], image_path: Path) -> None:
    shape = get_shape(slide, shape_path)
    image_part, r_id = shape.part.get_or_add_image_part(str(image_path))
    shape._pic.blipFill.blip.rEmbed = r_id

    with Image.open(image_path) as img:
        img_ratio = img.width / img.height
    box_ratio = shape.width / shape.height

    if img_ratio > box_ratio:
        target_w = img.height * box_ratio
        crop = (img.width - target_w) / img.width / 2
        shape.crop_left = crop
        shape.crop_right = crop
        shape.crop_top = 0
        shape.crop_bottom = 0
    else:
        target_h = img.width / box_ratio
        crop = (img.height - target_h) / img.height / 2
        shape.crop_top = crop
        shape.crop_bottom = crop
        shape.crop_left = 0
        shape.crop_right = 0


def clear_notes(prs: Presentation) -> None:
    for slide in prs.slides:
        try:
            tf = slide.notes_slide.notes_text_frame
        except Exception:
            continue
        if tf is None:
            continue
        for paragraph in tf.paragraphs:
            if paragraph.runs:
                paragraph.runs[0].text = ""
                for run in paragraph.runs[1:]:
                    run.text = ""


def build() -> None:
    assets = ensure_assets()
    prs = Presentation(str(TEMPLATE))
    for slide in prs.slides:
        slide._element.attrib.pop("show", None)

    slide = prs.slides[0]
    set_paragraphs(slide.shapes[5], ["论文图解"])
    set_paragraphs(slide.shapes[7], ["LFM2-VL"])
    replace_picture(slide, [2], assets["cover_banner"])

    slide = prs.slides[1]
    set_paragraphs(slide.shapes[1], ["目录"])
    set_paragraphs(
        slide.shapes[4],
        [
            "• 为什么值得关注：边缘级 VLM 的部署挑战与 LFM2-VL 的定位。",
            "• 核心创新：原生分辨率、动态切片、轻量连接器与可调 token 预算。",
            "• 逐图解读：图1 Demo、图2 架构、图3 速度、图4 显存。",
            "• 评测与选型：benchmark 结果如何看，450M 与 1.6B 如何取舍。",
            "• 结论：这条路线对端侧多模态系统意味着什么。",
        ],
    )

    slide = prs.slides[2]
    set_paragraphs(slide.shapes[1], ["引言：为什么要关注边缘级视觉语言模型"])
    set_paragraphs(slide.shapes[2], ["真正难的不是做出 VLM，而是在端侧把它跑得又快又稳。"])
    set_paragraphs(
        get_shape(slide, [4, 3]),
        [
            "• 低延迟：拍照问答、现场巡检、AR 交互都要求即时响应。",
            "• 低资源：手机、笔记本和边缘盒子对内存与功耗都很敏感。",
            "• 灵活推理：不同图像大小和任务复杂度，不该被固定输入流程绑死。",
            "• LFM2-VL 的价值就在于：把“能跑起来”也当成模型设计目标。",
        ],
    )
    replace_picture(slide, [6], assets["edge_intro"])

    slide = prs.slides[3]
    set_paragraphs(slide.shapes[1], ["核心创新：把端侧 VLM 做成可调节的 token 经济学"])
    set_paragraphs(slide.shapes[2], ["原生分辨率、动态切片和轻量连接器共同服务部署效率。"])
    set_paragraphs(
        get_shape(slide, [4, 3]),
        [
            "• 基于 LFM2 语言主干扩展到视觉语言，而不是重新堆一个更重的骨干。",
            "• 图像原生处理到 512x512，避免为了统一输入而强行上采样。",
            "• 大图按 512x512 patch 切片，并用缩略图保留整体场景语义。",
            "• PixelUnshuffle + 2-layer MLP 先压缩视觉 token，再送入语言模型。",
        ],
    )
    replace_picture(slide, [6], assets["architecture"])

    slide = prs.slides[4]
    set_paragraphs(slide.shapes[1], ["先补背景：LFM2 给 LFM2-VL 提供了什么底座"])
    set_paragraphs(slide.shapes[2], ["没看过 LFM2 也只要记住，它本身就是面向效率优化的语言模型。"])
    set_paragraphs(
        slide.shapes[5],
        [
            "• LFM2 主干强调短卷积与少量 GQA 的混合设计，目标是兼顾表达能力与推理速度。",
            "• 这意味着 VLM 部分继承的是一个本来就面向端侧和低延迟优化过的语言 backbone。",
            "• LFM2-VL-450M 使用约 86M 的 SigLIP2 NaFlex base 编码器，优先考虑更紧的资源预算。",
            "• LFM2-VL-1.6B 使用约 400M 的 shape-optimized 编码器，换来更强的细粒度视觉能力。",
            "• 两个版本共享同一套方法论：先把成本控住，再尽量把视觉信息保留下来。",
        ],
    )

    slide = prs.slides[5]
    set_paragraphs(slide.shapes[1], ["逐图导览：这篇工作主要靠哪几张图说服我们"])
    set_paragraphs(slide.shapes[2], ["四张图分别回答：能做什么、怎么做、快多少、省多少。"])
    set_paragraphs(
        slide.shapes[5],
        [
            "• 图1：用一个直观 demo 展示 450M 模型的图像理解与语言描述能力。",
            "• 图2：解释视觉编码、token 压缩和 LFM2 主干如何拼成完整的 VLM。",
            "• 表1：把精度放回参数量和部署成本语境里，观察它是否“对得起体量”。",
            "• 图3：速度对比证明其 GPU 推理时间能明显压低到更适合交互应用的区间。",
            "• 图4：显存对比说明它不是绝对最小，但在效率-能力间找到了更平衡的位置。",
        ],
    )

    slide = prs.slides[6]
    set_paragraphs(slide.shapes[1], ["图1解读：一个小 demo 实际上在展示什么"])
    set_paragraphs(slide.shapes[2], ["它不是 benchmark，但很直观地展示了 450M 版本的基本能力边界。"])
    set_paragraphs(
        get_shape(slide, [4, 3]),
        [
            "• 模型先识别出海滩场景，再抓住“牛在海边”这种相对少见的关键细节。",
            "• 输出不是单标签分类，而是完整自然语言描述，说明视觉 grounding 与生成是连通的。",
            "• 论文故意选 450M 版本做展示，强调的正是“小模型也能完成有用视觉理解”。",
            "• 当然，这类 demo 更像能力切片，真正可靠性还要回到后面的 benchmark 与效率图来看。",
        ],
    )
    replace_picture(slide, [6], assets["demo_showcase"])

    slide = prs.slides[7]
    set_paragraphs(slide.shapes[1], ["图2解读：LFM2-VL 的整体架构是怎样工作的"])
    set_paragraphs(slide.shapes[2], ["核心不是把图像硬塞进 LLM，而是先把视觉信息结构化压缩。"])
    set_paragraphs(
        get_shape(slide, [4, 3]),
        [
            "• 底层语言推理由 LFM2 backbone 承担，视觉部分只负责把图像变成可对齐的 token。",
            "",
            "• SigLIP2 NaFlex 编码器既支持原生分辨率输入，也支持非标准长宽比，不必先统一缩放。",
            "",
            "• 中间的 connector 通过 PixelUnshuffle + MLP 降低 token 数量，再和文本 token 共同进入 LFM2。",
            "",
            "• 这让模型的“重活”更多集中在有效信息而非冗余像素上，是效率优势的关键来源。",
        ],
    )
    replace_picture(slide, [6], assets["architecture"])

    slide = prs.slides[8]
    set_paragraphs(slide.shapes[1], ["图2继续：高分辨率为什么还能跑得动"])
    set_paragraphs(slide.shapes[2], ["关键在动态切片、全局缩略图，以及可在线调节的视觉 token 预算。"])
    set_paragraphs(
        get_shape(slide, [4, 3]),
        [
            "• 小图直接按原生尺寸处理，避免先放大再编码带来的额外开销。",
            "• 大图会被切成多个 512x512 patch；1.6B 版本还会额外看到一个全局 thumbnail。",
            "• 论文给出的例子是：256x384 约 96 tokens，384x680 约 240 tokens，1000x3000 约 1020 tokens。",
            "• 用户在推理时可以直接调最大图像 token 数和 patch 数量，不必重新训练模型。",
        ],
    )
    replace_picture(slide, [6], assets["token_flow"])

    slide = prs.slides[9]
    set_paragraphs(slide.shapes[1], ["评测结果：不是绝对 SOTA，但非常对得起体量"])
    set_paragraphs(slide.shapes[2], ["要把精度放回参数规模与部署成本语境里，才能看懂这篇论文的卖点。"])
    set_paragraphs(
        get_shape(slide, [4, 3]),
        [
            "• 1.6B 版本在 RealWorldQA 上到 65.23，已经接近 2B 级 InternVL3-2B 的 65.10。",
            "• OCRBench、MathVista、SEEDBench_IMG 等结果说明它在 OCR、高分辨率和推理类任务上并不掉队。",
            "• 450M 版本不是拿来冲榜的，而是给手机、穿戴和更苛刻设备预算准备的可用选项。",
            "• 所以论文真正强调的是“足够好的精度 + 明显更好的效率”，而不是一味追求最高分。",
        ],
    )
    replace_picture(slide, [6], assets["benchmark"])

    slide = prs.slides[10]
    set_paragraphs(slide.shapes[1], ["图3解读：速度优势到底有多实在"])
    set_paragraphs(slide.shapes[2], ["在相同任务设定下，它把交互延迟直接压到了更适合产品落地的区间。"])
    set_paragraphs(
        get_shape(slide, [4, 3]),
        [
            "• 论文设定的 workload 是：1 张 1024x1024 图像 + 简短提示词 + 生成 100 个 token。",
            "• LFM2-VL-1.6B 的处理时间约 833ms，而几款对比模型分别在 1650ms、1837ms 和 1950ms 左右。",
            "• 这意味着它相对最快竞品也接近 2 倍速度提升，优势不是小修小补，而是交互层面的体感差异。",
            "• 这张图最能说明：token 压缩和高效 backbone 的设计，最终都会兑现成真实等待时间的减少。",
        ],
    )
    replace_picture(slide, [6], assets["speed"])

    slide = prs.slides[11]
    set_paragraphs(slide.shapes[1], ["图4解读：显存并非最小，但性价比很高"])
    set_paragraphs(slide.shapes[2], ["要理解这张图，重点不是“谁最低”，而是“谁在更强能力下仍没把成本拉爆”。"])
    set_paragraphs(
        get_shape(slide, [4, 3]),
        [
            "• LFM2-VL-1.6B 约 3.14GB，低于 2B 级对手的 4.20GB 和 4.62GB，但高于更小的 1B 模型 2.06GB。",
            "• 也就是说，它不是绝对最省内存的模型，但在参数规模、速度和视觉能力之间更平衡。",
            "• 对笔记本、本地工作站和单 GPU 场景，这种“稍高一点显存换来更强能力”的交易通常是划算的。",
            "• 反过来，如果预算极紧，450M 版本才是更符合边缘设备约束的选择。",
        ],
    )
    replace_picture(slide, [6], assets["memory"])

    slide = prs.slides[12]
    set_paragraphs(slide.shapes[1], ["部署决策：450M 和 1.6B 应该怎么选"])
    set_paragraphs(slide.shapes[4], ["先按设备预算分层，再看任务是否依赖细粒度视觉理解。"])
    set_paragraphs(
        slide.shapes[6],
        [
            "450M 搭载更轻的视觉编码器，更适合手机、穿戴、边缘盒子这类强调实时性与低资源占用的场景；1.6B 则更适合笔记本、单 GPU 和需要更强 OCR / 高分辨率理解的任务。",
        ],
    )
    set_paragraphs(
        slide.shapes[7],
        [
            "两者共享同一套 token 预算调节机制，所以推荐的选型顺序是：先看内存与延迟，再判断任务是否依赖细节，最后再调图像 token 上限和 patch 数量。",
        ],
    )
    replace_picture(slide, [6], assets["variant_compare"])

    slide = prs.slides[13]
    set_paragraphs(slide.shapes[1], ["综合讨论：这篇工作真正贡献了什么"])
    set_paragraphs(slide.shapes[2], ["它把端侧 VLM 的优化目标说清楚了：不是更大，而是更可部署。"])
    set_paragraphs(
        slide.shapes[5],
        [
            "• 这篇论文最有价值的地方，不是某个单点分数，而是把“速度、显存、精度”一起纳入设计目标。",
            "• LFM2-VL 说明高分辨率视觉理解不一定等于固定高成本，token 预算本身可以成为可调控制杆。",
            "• 对工程侧来说，这比单纯追求 benchmark 榜首更有意义，因为它直接影响产品是否能落地。",
            "• 真正的启发是：未来端侧多模态系统的竞争，很可能是体系化效率设计的竞争。",
        ],
    )

    slide = prs.slides[14]
    set_paragraphs(slide.shapes[1], ["总结与结语"])
    set_paragraphs(slide.shapes[2], ["LFM2-VL 代表了一条更务实的边缘多模态路线。"])
    set_paragraphs(
        slide.shapes[5],
        [
            "• 图1告诉我们：即便是 450M 级别，小模型也能完成有用的视觉 grounding 与描述。",
            "• 图2给出的配方是全篇核心：NaFlex 编码器、动态切片、thumbnail 和 token 压缩共同协作。",
            "• 图3和图4进一步证明：效率收益是真实存在的，而且能直接转化成更短等待时间与更低部署门槛。",
            "• 如果目标是做边缘侧多模态系统，LFM2-VL 的最大启示是“可控成本”往往比 headline SOTA 更重要。",
        ],
    )

    clear_notes(prs)
    prs.save(str(OUTPUT))


if __name__ == "__main__":
    build()
