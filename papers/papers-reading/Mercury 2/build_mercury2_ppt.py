from __future__ import annotations

from pathlib import Path
from shutil import copyfile

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.util import Pt


BASE = Path("/Users/bingranyou/Library/Mobile Documents/com~apple~CloudDocs/Downloads/huaying")
WORKDIR = BASE / "Mercury 2"
TEMPLATE = BASE / "PPT模板.pptx"
OUTPUT = WORKDIR / "Mercury 2_扩散推理速览_中文版.pptx"
ASSET_DIR = WORKDIR / "_ppt_assets"


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def fit_canvas_size(ratio: float, long_side: int = 2400) -> tuple[int, int]:
    if ratio >= 1:
        width = long_side
        height = max(800, int(round(width / ratio)))
    else:
        height = long_side
        width = max(800, int(round(height * ratio)))
    return width, height


def prepare_image(src: Path, name: str, ratio: float, mode: str = "contain", background: str = "#ffffff") -> Path:
    ensure_dir(ASSET_DIR)
    width, height = fit_canvas_size(ratio)
    dst = ASSET_DIR / name
    with Image.open(src) as im:
        im = ImageOps.exif_transpose(im).convert("RGB")
        if mode == "cover":
            prepared = ImageOps.fit(im, (width, height), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))
        else:
            prepared = Image.new("RGB", (width, height), background)
            resized = ImageOps.contain(im, (width, height), method=Image.Resampling.LANCZOS)
            x = (width - resized.width) // 2
            y = (height - resized.height) // 2
            prepared.paste(resized, (x, y))
        prepared.save(dst, quality=95)
    return dst


def stack_vertical(
    sources: list[Path],
    name: str,
    ratio: float,
    background: str = "#ffffff",
    gap: int = 36,
) -> Path:
    ensure_dir(ASSET_DIR)
    width, height = fit_canvas_size(ratio)
    dst = ASSET_DIR / name
    canvas = Image.new("RGB", (width, height), background)
    inner_h = height - gap * (len(sources) - 1)
    panel_h = inner_h // len(sources)
    y = 0
    for src in sources:
        with Image.open(src) as im:
            im = ImageOps.exif_transpose(im).convert("RGB")
            resized = ImageOps.contain(im, (width, panel_h), method=Image.Resampling.LANCZOS)
            x = (width - resized.width) // 2
            paste_y = y + (panel_h - resized.height) // 2
            canvas.paste(resized, (x, paste_y))
        y += panel_h + gap
    canvas.save(dst, quality=95)
    return dst


def stack_horizontal(
    sources: list[Path],
    name: str,
    ratio: float,
    background: str = "#ffffff",
    gap: int = 28,
) -> Path:
    ensure_dir(ASSET_DIR)
    width, height = fit_canvas_size(ratio)
    dst = ASSET_DIR / name
    canvas = Image.new("RGB", (width, height), background)
    inner_w = width - gap * (len(sources) - 1)
    panel_w = inner_w // len(sources)
    x = 0
    for src in sources:
        with Image.open(src) as im:
            im = ImageOps.exif_transpose(im).convert("RGB")
            resized = ImageOps.contain(im, (panel_w, height), method=Image.Resampling.LANCZOS)
            y = (height - resized.height) // 2
            paste_x = x + (panel_w - resized.width) // 2
            canvas.paste(resized, (paste_x, y))
        x += panel_w + gap
    canvas.save(dst, quality=95)
    return dst


def replace_picture(slide, shape_index: int, image_path: Path) -> None:
    old = slide.shapes[shape_index - 1]
    left, top, width, height = old.left, old.top, old.width, old.height
    parent = old._element.getparent()
    insert_at = list(parent).index(old._element)
    parent.remove(old._element)
    new_pic = slide.shapes.add_picture(str(image_path), left, top, width=width, height=height)
    parent.remove(new_pic._element)
    parent.insert(insert_at, new_pic._element)


def set_plain_text(shape, text: str, size: float | None = None, bold: bool | None = None) -> None:
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    for run in p.runs:
        if size is not None:
            run.font.size = Pt(size)
        if bold is not None:
            run.font.bold = bold


def set_bullets(shape, items: list[str], size: float = 13, bold_prefix: bool = False) -> None:
    tf = shape.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(size)
            if bold_prefix:
                run.font.bold = True


def set_group_bullets(slide, group_shape_index: int, items: list[str], size: float = 13) -> None:
    group = slide.shapes[group_shape_index - 1]
    text_shape = group.shapes[1]
    set_bullets(text_shape, items, size=size)


def build_assets() -> dict[str, Path]:
    intro_dir = WORKDIR / "Introducing Mercury 2 – Inception_files"
    zh_dir = WORKDIR / "扩散模型成最快深度思考！告别自回归每秒1009个tokens，英伟达微软都投了_files"
    assets = {
        "cover": prepare_image(zh_dir / "640(9)", "cover.png", ratio=12.04 / 3.95, mode="cover"),
        "parallel": prepare_image(zh_dir / "640", "parallel.png", ratio=9.74 / 3.67, mode="contain"),
        "speed_benchmark": prepare_image(intro_dir / "KUQ2ijMteh6pTFCieisdOx9s.png", "speed_benchmark.png", ratio=7.84 / 5.76, mode="contain"),
        "official_table": prepare_image(intro_dir / "vUhManfsMjNnhb8We1IfhRGbfr8.png", "official_table.png", ratio=9.74 / 3.67, mode="contain"),
        "intelligence_speed": prepare_image(zh_dir / "640(5)", "intelligence_speed.png", ratio=7.84 / 5.76, mode="contain"),
        "speed_price": prepare_image(zh_dir / "640(8)", "speed_price.png", ratio=10.37 / 3.91, mode="contain"),
        "arxiv": prepare_image(zh_dir / "640(10)", "arxiv.png", ratio=9.48 / 3.57, mode="contain"),
        "market_reaction": prepare_image(WORKDIR / "market_reaction_collage.png", "market_reaction.png", ratio=6.12 / 5.98, mode="contain"),
        "bar_speed": prepare_image(zh_dir / "640(6)", "bar_speed.png", ratio=5.81 / 2.09, mode="contain"),
    }
    assets["team_combo"] = stack_horizontal(
        [intro_dir / "RoZiv6hJStkiJmnDuDeCk9S2s.webp", zh_dir / "640(11)"],
        "team_combo.png",
        ratio=8.95 / 3.37,
        background="#ffffff",
    )
    assets["metrics_stack"] = stack_vertical(
        [intro_dir / "KUQ2ijMteh6pTFCieisdOx9s.png", intro_dir / "vUhManfsMjNnhb8We1IfhRGbfr8.png"],
        "metrics_stack.png",
        ratio=5.81 / 2.09,
        background="#ffffff",
        gap=18,
    )
    return assets


def build_presentation() -> None:
    ensure_dir(ASSET_DIR)
    assets = build_assets()
    copyfile(TEMPLATE, OUTPUT)
    prs = Presentation(OUTPUT)
    for slide in prs.slides:
        slide._element.attrib.pop("show", None)

    # Slide 1
    slide = prs.slides[0]
    set_plain_text(slide.shapes[5], "2026年3月", size=16, bold=False)
    set_plain_text(slide.shapes[7], "Mercury 2\n扩散推理速览", size=28, bold=True)
    replace_picture(slide, 2, assets["cover"])

    # Slide 2
    slide = prs.slides[1]
    set_plain_text(slide.shapes[1], "目录", size=24, bold=True)
    set_bullets(
        slide.shapes[3],
        [
            "Mercury 2 的模型定位与核心亮点",
            "扩散式并行生成为何能把速度拉到 1009 tokens/s",
            "性能、价格与应用场景的综合判断",
            "Inception Labs 的技术路线、团队与产业背书",
            "对实时 AI 产品与投资观察的启发",
        ],
        size=18,
    )

    # Slide 3
    slide = prs.slides[2]
    set_plain_text(slide.shapes[1], "模型定位：Mercury 2 是什么", size=22, bold=True)
    set_plain_text(slide.shapes[2], "全球最快的推理型 LLM 之一，核心卖点是“扩散式并行生成”", size=13, bold=False)
    set_group_bullets(
        slide,
        4,
        [
            "官方将其定义为推理型语言模型，但底层不再沿用逐 token 自回归解码。",
            "生成方式更像“先形成草稿、再并行修订”的编辑过程，可同时优化多个 token。",
            "对于 Agent、检索、抽取、代码补全等多步循环任务，延迟不会像传统模型那样线性累积。",
        ],
        size=13,
    )
    replace_picture(slide, 6, assets["parallel"])

    # Slide 4
    slide = prs.slides[3]
    set_plain_text(slide.shapes[1], "为什么它会快：速度曲线被重写", size=22, bold=True)
    set_plain_text(slide.shapes[2], "从“打字机模式”切换到并行 refinement，是 Mercury 2 的根本差异", size=13, bold=False)
    set_group_bullets(
        slide,
        4,
        [
            "官方披露在 NVIDIA Blackwell GPU 上可达 1009 tokens/s。",
            "相较 GPT-5 mini、Claude 4.5 Haiku 等轻量模型，速度优势超过 5 倍。",
            "快的不只是首 token，更关键是整段输出吞吐更高，适合高并发生产环境。",
            "第三方传播里，Mercury 2 已经被视为“扩散式推理”进入主流视野的标志事件。",
        ],
        size=12.5,
    )
    replace_picture(slide, 6, assets["speed_benchmark"])

    # Slide 5
    slide = prs.slides[4]
    set_plain_text(slide.shapes[1], "指标总览：速度、价格、能力一起看", size=22, bold=True)
    set_plain_text(slide.shapes[2], "Mercury 2 不是单纯“快模型”，而是在重画质量-速度-成本曲线", size=13, bold=False)
    set_bullets(
        slide.shapes[4],
        [
            "速度：1009 tokens/s，主打实时推理与高并发响应。",
            "价格：输入 0.25 美元 / 百万，输出 0.75 美元 / 百万。",
            "能力：GPQA、LCB、AIME 等基准不输主流轻量/速度型模型。",
            "功能：128K 上下文、工具调用、结构化 JSON、可调推理。",
            "接口：兼容 OpenAI API，迁移与接入成本较低。",
        ],
        size=17,
    )

    # Slide 6
    slide = prs.slides[5]
    set_plain_text(slide.shapes[1], "应用场景：速度优势如何转化为价值", size=22, bold=True)
    set_plain_text(slide.shapes[2], "真正受益的是“多轮、多步、强交互”的生产型 AI 工作流", size=13, bold=False)
    set_bullets(
        slide.shapes[4],
        [
            "代码与编辑：自动补全、改写、重构等场景对延迟极敏感。",
            "Agent：单任务常串联数十次推理，单次提速会被放大成整体收益。",
            "语音：让“推理质量”更接近自然对话的实时节奏。",
            "搜索 / RAG：多跳检索、重排与总结链路更易压到可接受延迟。",
            "企业部署：官方强调 p95 延迟、高并发稳定性和轮次一致性。",
        ],
        size=17,
    )

    # Slide 7
    slide = prs.slides[6]
    set_plain_text(slide.shapes[1], "能力与延迟表现：并没有为了速度显著降智", size=22, bold=True)
    set_plain_text(slide.shapes[2], "官方对比图显示，Mercury 2 在关键基准上接近甚至超过多款速度型模型", size=13, bold=False)
    set_group_bullets(
        slide,
        4,
        [
            "E2E latency 约 1.7 秒，显著低于表中多数对比模型。",
            "GPQA 74、LCB 67、AIME 91，说明其在科学问答、代码和数学上保持了较强竞争力。",
            "第三方中文解读特别强调：Mercury 2 在 AIME 上甚至压过部分以性能见长的 Flash 系列模型。",
        ],
        size=13,
    )
    replace_picture(slide, 6, assets["official_table"])

    # Slide 8
    slide = prs.slides[7]
    set_plain_text(slide.shapes[1], "第三方视角：Mercury 2 位于高速度高智能象限", size=22, bold=True)
    set_plain_text(slide.shapes[2], "若第三方基准可信，它已进入“小模型最优象限”", size=13, bold=False)
    set_group_bullets(
        slide,
        4,
        [
            "Artificial Analysis 图中，Mercury 2 被单独甩到 1k+ tokens/s 区间。",
            "在相近成本带内，速度与智能度组合都明显优于传统自回归速度型模型。",
            "这更像“架构范式红利”，而不是简单的工程加速。",
            "因此市场讨论的焦点，已从“能不能用”转向“会不会改写实时 AI 产品设计”。",
        ],
        size=12.5,
    )
    replace_picture(slide, 6, assets["intelligence_speed"])

    # Slide 9
    slide = prs.slides[8]
    set_plain_text(slide.shapes[1], "价格与产能：单位成本不高，吞吐优势更强", size=22, bold=True)
    set_plain_text(slide.shapes[2], "对生产部署而言，性价比来自“更低等待时间 × 更高吞吐”", size=13, bold=False)
    set_group_bullets(
        slide,
        4,
        [
            "在第三方图里，Mercury 2 的速度表现呈明显断层领先。",
            "价格并未因为追求极速而出现极端溢价，输入和输出单价仍在可部署区间。",
            "对于并发执行的 Agent、搜索、客服与语音场景，吞吐优势会进一步摊薄系统成本。",
            "这也是官方反复强调“production AI should feel instant”的底层原因。",
        ],
        size=12.5,
    )
    replace_picture(slide, 6, assets["speed_price"])

    # Slide 10
    slide = prs.slides[9]
    set_plain_text(slide.shapes[1], "技术根基：离散扩散让文本生成脱离自回归", size=22, bold=True)
    set_plain_text(slide.shapes[2], "Mercury 2 背后不是“优化版 Transformer”，而是一条持续多年的扩散路线", size=13, bold=False)
    set_group_bullets(
        slide,
        4,
        [
            "Stefano Ermon 团队长期研究如何把扩散模型从连续空间推广到离散 token 空间。",
            "2023 年的 SEDD 论文通过 score entropy 损失，把离散扩散真正带入文本建模。",
            "量子位文章提到，该论文后来获得 ICML 2024 最佳论文，为 Mercury 系列商业化提供了学术地基。",
            "Mercury 2 可以被视为这条学术路线向商业系统化落地的代表产品。",
        ],
        size=12.2,
    )
    replace_picture(slide, 6, assets["arxiv"])

    # Slide 11
    slide = prs.slides[10]
    set_plain_text(slide.shapes[1], "公司与团队：Inception Labs 持续押注扩散路线", size=22, bold=True)
    set_plain_text(slide.shapes[2], "不是跟风切换路线，而是围绕“速度与成本曲线”做长期下注", size=13, bold=False)
    set_group_bullets(
        slide,
        4,
        [
            "公司成立于 2024 年，使命是用 diffusion 机制替代传统自回归生成。",
            "CEO Stefano Ermon 出身斯坦福；核心团队还包括 UCLA 的 Aditya Grover 与 Cornell 的 Volodymyr Kuleshov。",
            "从官方文章到媒体解读都能看出，这家公司卖点非常聚焦：让“生产 AI 感觉像实时系统”。",
            "Mercury 2 的发布，本质上是在向市场证明扩散式语言模型具备商业可交付性。",
        ],
        size=12.2,
    )
    replace_picture(slide, 6, assets["team_combo"])

    # Slide 12
    slide = prs.slides[11]
    set_plain_text(slide.shapes[1], "产业背书与市场反响：英伟达、微软系与开发者都在关注", size=22, bold=True)
    set_plain_text(slide.shapes[2], "Mercury 2 的传播并不只是技术圈自嗨，而是直接连接到基础设施与应用落地", size=13, bold=False)
    set_group_bullets(
        slide,
        4,
        [
            "融资层面，量子位文章提到公司曾获 5000 万美元融资，投资方包括 NVentures、M12、Menlo Ventures 等。",
            "NVIDIA 官方账号第一时间祝贺，强调其在 Blackwell GPU 上的扩散推理速度。",
            "开发者讨论集中在端到端语音、实时 HCI、广告优化和自动化执行等对延迟极敏感的场景。",
            "这说明 Mercury 2 的价值更可能先在“实时系统”而非“离线 benchmark”里兑现。",
        ],
        size=12.2,
    )
    replace_picture(slide, 6, assets["market_reaction"])

    # Slide 13
    slide = prs.slides[12]
    set_plain_text(slide.shapes[1], "产品就绪度：已经具备进入企业工作流的条件", size=22, bold=True)
    set_plain_text(slide.shapes[4], "接口兼容：官方明确支持 OpenAI API 兼容格式，可直接嵌入既有调用链。", size=14, bold=False)
    set_plain_text(slide.shapes[6], "功能形态：128K 上下文、原生 tool use、可调 reasoning、schema-aligned JSON 输出。", size=12.5, bold=False)
    set_plain_text(slide.shapes[7], "部署含义：对于希望压缩 p95 延迟、提升并发吞吐、减少多步链路等待的团队，Mercury 2 已经不是“概念验证”，而是可直接试点评估的产品。", size=12.5, bold=False)
    replace_picture(slide, 6, assets["bar_speed"])

    # Slide 14
    slide = prs.slides[13]
    set_plain_text(slide.shapes[1], "投资与产品观察：Mercury 2 值得关注什么", size=22, bold=True)
    set_plain_text(slide.shapes[2], "它未必立刻替代所有模型，但很可能改写“实时 AI”这条产品线", size=13, bold=False)
    set_bullets(
        slide.shapes[4],
        [
            "范式变量最关键：若扩散式语言模型持续成立，自回归速度护城河会被改写。",
            "最先受益的不是所有任务，而是强交互、强并发、长链路的实时 AI 产品。",
            "短期风险在于生态成熟度、复杂工具调用稳定性与公开评测覆盖度仍待补强。",
            "若后续开放更多客户案例、成本曲线和大规模服务数据，Mercury 系列有机会形成差异化赛道。",
        ],
        size=17,
    )

    # Slide 15
    slide = prs.slides[14]
    set_plain_text(slide.shapes[1], "总结与结论", size=22, bold=True)
    set_plain_text(slide.shapes[2], "Mercury 2 展示了“扩散 + 推理 + 实时交互”结合后的新可能", size=13, bold=False)
    set_bullets(
        slide.shapes[4],
        [
            "核心亮点是快：1009 tokens/s 把速度从优化问题提升为架构问题。",
            "核心价值是稳：在多个 benchmark 上并未因极速而明显牺牲能力。",
            "核心壁垒是路线：公司、论文、融资与算力合作构成了完整叙事闭环。",
            "结论：Mercury 2 适合作为实时 Agent、语音与搜索工作流的重点跟踪标的。",
        ],
        size=17,
    )

    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build_presentation()
